"""Turning an uploaded file into words a voice can read.

Every document format lives here and nothing else does: no FastAPI, no
database, no settings object. Bytes and a filename in, text out.

Two things shape it. The output is spoken, not displayed, so layout fidelity is
worthless and sentence shape is everything — a spreadsheet read cell by cell is
noise, "Row 1. Region: North. Units: 412." is a sentence. And the files come
from strangers: Office and OpenDocument files are zips of XML, so nothing
reaches a parser before `_guard_zip` has looked at it.
"""

from __future__ import annotations

import csv
import io
import logging
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections.abc import Callable, Iterable, Iterator
from dataclasses import dataclass
from datetime import date, datetime
from html.parser import HTMLParser
from pathlib import Path

from .markdown_speech import DEFAULT_FORMAT

log = logging.getLogger("oneread.extract")

MARKDOWN = "markdown"

# Ceilings. A document past any of these is read up to the limit and marked
# truncated rather than refused, because half a book is more use than an error.
MAX_PAGES = 2_000
MAX_SLIDES = 1_000
MAX_SHEETS = 50
MAX_ROWS = 50_000
MAX_COLUMNS = 200

# Zip container guards.
MAX_ZIP_MEMBERS = 2_000
MAX_ZIP_RATIO = 200  # a real docx manages five- to ten-fold
ZIP_RATIO_FLOOR = 8 * 1024 * 1024  # below this, a high ratio is just a tidy file
ZIP_CHUNK = 256 * 1024
#: Carried between chunks so a marker straddling the boundary is still seen.
ENTITY_OVERLAP = 16

# A PDF with less text than this per page is pictures of pages, not a document.
MIN_PDF_CHARS_PER_PAGE = 20

SOFFICE_TIMEOUT_S = 120


class UnreadableFile(ValueError):
    """A file we can't turn into text.

    The message is shown to the person who uploaded it, word for word, so it
    has to be a sentence rather than a diagnostic.
    """


@dataclass(frozen=True)
class Extracted:
    text: str
    #: "plain" or "markdown" — what the entry's format becomes.
    format: str
    #: The document's own title where it has one, otherwise the filename.
    title: str
    #: Broad shape, for the wording in the interface: "document", "slides"…
    kind: str
    #: True when a ceiling was reached and the tail was left behind.
    truncated: bool = False


@dataclass(frozen=True)
class _Job:
    filename: str
    limit: int
    max_unzipped: int
    soffice: str


@dataclass(frozen=True)
class Reader:
    ext: str
    #: What to call this in the interface. Shown in the "what can I upload" list.
    label: str
    kind: str
    read: Callable[[bytes, _Builder, _Job], str]
    format: str = DEFAULT_FORMAT
    zipped: bool = False
    magic: bytes = b""
    #: Shown when the file parses but holds no words at all.
    empty: str = "There are no words in that file."


# --- collecting text ---------------------------------------------------------


class _Builder:
    """Gathers blocks of speech and stops once the character limit is passed.

    Readers push blocks in as they walk the document. Checking `full` between
    pages means a thousand-page PDF is abandoned early instead of assembled in
    memory and then thrown away.
    """

    def __init__(self, limit: int) -> None:
        self.limit = max(1, limit)
        self.parts: list[str] = []
        self.size = 0
        self.truncated = False

    @property
    def full(self) -> bool:
        return self.size >= self.limit

    def add(self, text: str) -> None:
        text = _tidy(text)
        if not text:
            return
        if self.full:
            self.truncated = True
            return
        self.parts.append(text)
        self.size += len(text) + 2

    def stop(self) -> None:
        """A ceiling other than the character limit was reached."""
        self.truncated = True

    def done(self) -> str:
        text = "\n\n".join(self.parts)
        if len(text) <= self.limit:
            return text

        self.truncated = True
        text = text[: self.limit]
        # Cutting mid-word reads badly, so fall back to the last sentence end
        # if there's one close enough to the cut to be worth keeping.
        edge = max(text.rfind(". "), text.rfind("\n"))
        return text[: edge + 1] if edge > self.limit - 400 else text


_SPACE = re.compile(r"[ \t ]+")
_BLANKS = re.compile(r"\n{3,}")


def _tidy(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _SPACE.sub(" ", text)
    return _BLANKS.sub("\n\n", text).strip()


def _sentence(text: str) -> str:
    """Give a fragment a full stop so the voice pauses after it.

    Headings, slide titles and cells aren't written as sentences, but they are
    read as ones, and without the stop they run into whatever follows.
    """
    text = _tidy(text)
    if not text:
        return ""
    return text if text[-1] in ".!?:;,…" else f"{text}."


# --- grids -------------------------------------------------------------------


def _cell(value: object) -> str:
    """A cell as a person would say it out loud."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "yes" if value else "no"
    if isinstance(value, datetime):
        said = f"{value.day} {value:%B %Y}"
        return said if value.time() == datetime.min.time() else f"{said} at {value:%H:%M}"
    if isinstance(value, date):
        return f"{value.day} {value:%B %Y}"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _looks_like_header(row: list[str]) -> bool:
    """Labels across the top, or just the first row of data?"""
    filled = [cell for cell in row if cell]
    if len(filled) < 2 or len(filled) < len(row) - 1:
        return False
    return not any(re.fullmatch(r"-?[\d.,%$€£]+", cell) for cell in filled)


def _grid(rows: Iterable[list[object]], build: _Builder, *, name: str = "") -> None:
    """Read a grid as sentences.

    Joining cells with commas turns a spreadsheet into forty minutes of numbers
    with no idea which column they came from. Naming the column before each
    value costs a few words and makes the recording usable.
    """
    header: list[str] | None = None
    count = 0
    said_name = False

    for raw in rows:
        if build.full:
            build.stop()
            return
        cells = [_cell(value) for value in list(raw)[:MAX_COLUMNS]]
        while cells and not cells[-1]:
            cells.pop()
        if not any(cells):
            continue

        if not said_name:
            if name:
                build.add(_sentence(f"Sheet: {name}"))
            said_name = True

        if header is None and _looks_like_header(cells):
            header = cells
            build.add(_sentence(f"Columns: {', '.join(cell for cell in cells if cell)}"))
            continue

        count += 1
        if count > MAX_ROWS:
            build.stop()
            return

        parts = [f"Row {count}."]
        for index, cell in enumerate(cells):
            if not cell:
                continue
            label = header[index] if header and index < len(header) and header[index] else ""
            parts.append(_sentence(f"{label}: {cell}" if label else cell))
        build.add(" ".join(parts))


# --- decoding ----------------------------------------------------------------


def _decode(data: bytes) -> str:
    """Text out of bytes, without a character-detection library.

    A byte-order mark settles it outright, and past that UTF-8 either decodes
    or it doesn't. Windows-1252 catches the older Western files, and the last
    attempt cannot raise, so this never fails on a file a person can open.
    """
    for bom, encoding in (
        (b"\xef\xbb\xbf", "utf-8-sig"),
        (b"\xff\xfe\x00\x00", "utf-32"),
        (b"\x00\x00\xfe\xff", "utf-32"),
        (b"\xff\xfe", "utf-16"),
        (b"\xfe\xff", "utf-16"),
    ):
        if data.startswith(bom):
            try:
                return data.decode(encoding)
            except UnicodeDecodeError:
                break

    for encoding in ("utf-8", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# --- zip safety --------------------------------------------------------------


TOO_BIG = "That file unpacks to far more than it looks like. Skipped."


def _oversized(unpacked: int, packed: int, max_unzipped: int) -> bool:
    if unpacked > max_unzipped:
        return True
    # Below the floor a high ratio is just a tidy file, not an attack.
    return unpacked > ZIP_RATIO_FLOOR and unpacked > packed * MAX_ZIP_RATIO


def _guard_zip(data: bytes, job: _Job) -> None:
    """Look over a zip container before any parser opens it.

    docx, pptx, xlsx and the OpenDocument formats are zips of XML, which brings
    two cheap attacks: an archive that expands to fill the disk, and an XML
    entity declaration that expands inside the parser.

    A zip's declared sizes are a claim, not a measurement, so they only serve as
    a fast reject; every member is then read in chunks and counted for real. One
    bounded pass, and a member outrunning its declared size is refused on the
    first chunk.

    The entity check is defence in depth — python-docx and python-pptx parse
    with `resolve_entities=False` and odfpy through `defusedxml.sax` — kept
    because it is nearly free once the bytes are going past anyway.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as bundle:
            members = bundle.infolist()
            if len(members) > MAX_ZIP_MEMBERS:
                raise UnreadableFile("There are too many parts inside that file to read it.")

            if _oversized(
                sum(member.file_size for member in members), len(data), job.max_unzipped
            ):
                raise UnreadableFile(TOO_BIG)

            total = 0
            for member in members:
                if member.is_dir():
                    continue
                scan = member.filename.lower().endswith((".xml", ".rels"))
                seen = 0
                tail = b""
                with bundle.open(member) as stream:
                    while True:
                        chunk = stream.read(ZIP_CHUNK)
                        if not chunk:
                            break
                        seen += len(chunk)
                        total += len(chunk)
                        if seen > member.file_size or _oversized(
                            total, len(data), job.max_unzipped
                        ):
                            raise UnreadableFile(TOO_BIG)
                        if scan:
                            window = tail + chunk
                            if b"<!DOCTYPE" in window or b"<!ENTITY" in window:
                                raise UnreadableFile(
                                    "That file has instructions inside it that "
                                    "oneread won't run."
                                )
                            tail = window[-ENTITY_OVERLAP:]
    except zipfile.BadZipFile:
        # A member whose real contents don't match its declared size fails its
        # checksum here, which is the same answer as a truncated file: we can't
        # read it, and we haven't spent anything finding that out.
        raise UnreadableFile(
            "That file is damaged, or it isn't the kind of file its name says it is."
        ) from None


# --- readers -----------------------------------------------------------------


def _read_text(data: bytes, build: _Builder, job: _Job) -> str:
    build.add(_decode(data))
    return ""


def _read_markdown(data: bytes, build: _Builder, job: _Job) -> str:
    """Markdown goes in untouched — the app already knows how to speak it."""
    text = _decode(data).replace("\r\n", "\n").replace("\r", "\n").strip()
    build.add(text)

    # A leading "# Heading" is almost always the document's name.
    first = text.split("\n", 1)[0].strip()
    return first.lstrip("#").strip() if first.startswith("#") else ""


def _read_csv(data: bytes, build: _Builder, job: _Job) -> str:
    text = _decode(data)
    try:
        dialect: type[csv.Dialect] | csv.Dialect = csv.Sniffer().sniff(
            text[:8192], delimiters=",;\t|"
        )
    except csv.Error:
        # Sniffer gives up on one-column files and on ragged ones. The
        # extension is the next best guess, and a plain comma after that.
        class _Fallback(csv.excel):
            delimiter = "\t" if job.filename.lower().endswith(".tsv") else ","

        dialect = _Fallback

    _grid(csv.reader(io.StringIO(text), dialect), build)
    return ""


def _read_docx(data: bytes, build: _Builder, job: _Job) -> str:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(io.BytesIO(data))

    # Walking `document.paragraphs` is the obvious move and the wrong one: it
    # skips every table, so a report's numbers vanish and the prose around them
    # reads as if they were never there. The body's own children keep both in
    # the order they were written.
    for child in document.element.body.iterchildren():
        if build.full:
            build.stop()
            break
        tag = str(child.tag).rsplit("}", 1)[-1]
        if tag == "p":
            paragraph = Paragraph(child, document)
            style = (paragraph.style.name or "") if paragraph.style is not None else ""
            text = paragraph.text
            build.add(_sentence(text) if style.startswith(("Heading", "Title")) else text)
        elif tag == "tbl":
            table = Table(child, document)
            _grid(([cell.text for cell in row.cells] for row in table.rows), build)

    return document.core_properties.title or ""


def _read_pptx(data: bytes, build: _Builder, job: _Job) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))

    for number, slide in enumerate(presentation.slides, start=1):
        if build.full or number > MAX_SLIDES:
            build.stop()
            break

        title_shape = slide.shapes.title
        title = _tidy(title_shape.text) if title_shape is not None else ""
        build.add(_sentence(f"Slide {number}") + (f" {_sentence(title)}" if title else ""))

        # Iterating the shapes builds fresh wrappers, so identity says nothing
        # about whether this is the title we just read out. The id does.
        title_id = title_shape.shape_id if title_shape is not None else None

        for shape in _shapes(slide.shapes):
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            if getattr(shape, "has_table", False):
                table = shape.table
                _grid(([cell.text for cell in row.cells] for row in table.rows), build)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    build.add(_sentence(paragraph.text))

        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            notes = _tidy(frame.text) if frame is not None else ""
            if notes:
                build.add(f"Notes. {notes}")

    return presentation.core_properties.title or ""


def _shapes(shapes: Iterable) -> Iterator:
    """Flatten grouped shapes — text hides inside groups all the time."""
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _shapes(shape.shapes)
        else:
            yield shape


def _read_xlsx(data: bytes, build: _Builder, job: _Job) -> str:
    from openpyxl import load_workbook

    # read_only streams rows instead of building the whole grid in memory, and
    # data_only hands back the last calculated value rather than the formula,
    # which is what somebody listening actually wants to hear.
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        sheets = workbook.worksheets[:MAX_SHEETS]
        if len(workbook.worksheets) > MAX_SHEETS:
            build.stop()
        several = len(sheets) > 1
        for sheet in sheets:
            if build.full:
                build.stop()
                break
            _grid(sheet.iter_rows(values_only=True), build, name=sheet.title if several else "")
    finally:
        workbook.close()
    return ""


def _read_xls(data: bytes, build: _Builder, job: _Job) -> str:
    import xlrd

    try:
        workbook = xlrd.open_workbook(file_contents=data)
    except Exception:
        raise UnreadableFile(
            "That spreadsheet is in a format oneread can't open. Save it as .xlsx and try again."
        ) from None

    sheets = workbook.sheets()[:MAX_SHEETS]
    several = len(sheets) > 1
    for sheet in sheets:
        if build.full:
            build.stop()
            break
        rows = (sheet.row_values(index) for index in range(sheet.nrows))
        _grid(rows, build, name=sheet.name if several else "")
    return ""


_HYPHEN_BREAK = re.compile(r"(\w)-\n(\w)")
_LINE_BREAK = re.compile(r"(?<![\n])\n(?!\n)")


def _read_pdf(data: bytes, build: _Builder, job: _Job) -> str:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(data))
    except PdfReadError:
        raise UnreadableFile("That PDF is damaged and can't be opened.") from None

    if reader.is_encrypted:
        try:
            opened = reader.decrypt("")
        except Exception:
            opened = 0
        if not opened:
            raise UnreadableFile("That PDF is password-protected. Remove the password and retry.")

    pages = list(reader.pages)
    read = 0
    letters = 0

    for page in pages[:MAX_PAGES]:
        if build.full:
            build.stop()
            break
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        read += 1
        letters += len(text.strip())

        # PDFs carry no paragraphs, only lines that happened to end where the
        # page did. Rejoining them keeps sentences whole for the segmenter;
        # exact paragraph shape doesn't survive being read aloud anyway.
        text = _HYPHEN_BREAK.sub(r"\1\2", text)
        build.add(_LINE_BREAK.sub(" ", text))

    if len(pages) > MAX_PAGES:
        build.stop()

    if read and letters < max(40, MIN_PDF_CHARS_PER_PAGE * read):
        raise UnreadableFile(
            "That PDF is pictures of pages rather than text. oneread can't read those yet. "
            "A PDF you can select text in will work."
        )

    return (reader.metadata.title or "") if reader.metadata else ""


def _read_rtf(data: bytes, build: _Builder, job: _Job) -> str:
    from striprtf.striprtf import rtf_to_text

    build.add(rtf_to_text(_decode(data), errors="ignore"))
    return ""


class _Stripper(HTMLParser):
    """HTML down to the words, using what's already in the standard library."""

    SKIP = {"script", "style", "head", "noscript", "template", "svg"}
    BLOCK = {
        "p", "div", "section", "article", "br", "li", "tr", "h1", "h2", "h3",
        "h4", "h5", "h6", "blockquote", "pre", "td", "th",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.title = ""
        self._depth = 0
        self._in_title = False

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self.SKIP:
            self._depth += 1
        elif tag == "title":
            self._in_title = True
        elif tag in self.BLOCK:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self.SKIP and self._depth:
            self._depth -= 1
        elif tag == "title":
            self._in_title = False
        elif tag in self.BLOCK:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._in_title:
            self.title += data
        elif not self._depth:
            self.parts.append(data)


def _read_html(data: bytes, build: _Builder, job: _Job) -> str:
    stripper = _Stripper()
    stripper.feed(_decode(data))
    stripper.close()
    build.add("".join(stripper.parts))
    return stripper.title.strip()


def _read_odf(data: bytes, build: _Builder, job: _Job) -> str:
    from odf import table as odf_table
    from odf import teletype
    from odf.dc import Title
    from odf.draw import Page
    from odf.opendocument import load

    try:
        document = load(io.BytesIO(data))
    except Exception:
        raise UnreadableFile("That file is damaged, or it isn't an OpenDocument file.") from None

    body = document.body
    kind = str(document.mimetype or "")
    found = document.meta.getElementsByType(Title) if document.meta is not None else []
    title = teletype.extractText(found[0]) if found else ""

    if "presentation" in kind:
        for number, page in enumerate(body.getElementsByType(Page), start=1):
            if build.full or number > MAX_SLIDES:
                build.stop()
                break
            build.add(_sentence(f"Slide {number}"))
            for line in teletype.extractText(page).split("\n"):
                build.add(_sentence(line))
        return title

    if "spreadsheet" in kind:
        sheets = body.getElementsByType(odf_table.Table)[:MAX_SHEETS]
        several = len(sheets) > 1
        for sheet in sheets:
            if build.full:
                build.stop()
                break
            _grid(
                _odf_rows(sheet),
                build,
                name=str(sheet.getAttribute("name") or "") if several else "",
            )
        return title

    # `body` is office:body, whose one child holds the actual content, so the
    # paragraphs are a level further down than they look.
    for node in _odf_content(body):
        if build.full:
            build.stop()
            break
        tag = node.qname[1] if getattr(node, "qname", None) else ""
        if tag == "table":
            _grid(_odf_rows(node), build)
        elif tag in {"h", "p", "list"}:
            text = teletype.extractText(node)
            build.add(_sentence(text) if tag == "h" else text)

    return title


def _odf_content(body) -> Iterator:
    """The blocks inside office:body, whichever wrapper they came in."""
    for child in body.childNodes:
        tag = child.qname[1] if getattr(child, "qname", None) else ""
        if tag in {"text", "spreadsheet", "presentation", "drawing", "chart"}:
            yield from child.childNodes
        else:
            yield child


def _odf_rows(table) -> Iterator[list[str]]:
    """Rows out of an OpenDocument table, honouring repeat counts.

    A run of identical or empty cells is stored once with a repeat attribute.
    Ignoring it puts every value in the wrong column.
    """
    from odf import table as odf_table
    from odf import teletype

    for row in table.getElementsByType(odf_table.TableRow)[:MAX_ROWS]:
        cells: list[str] = []
        for cell in row.getElementsByType(odf_table.TableCell):
            value = teletype.extractText(cell)
            repeat = cell.getAttribute("numbercolumnsrepeated")
            times = min(int(repeat), MAX_COLUMNS) if repeat and str(repeat).isdigit() else 1
            if not value and times > 8:
                times = 1  # trailing padding to the sheet's edge, not real columns
            cells.extend([value] * times)
            if len(cells) >= MAX_COLUMNS:
                break
        yield cells


def _read_legacy(data: bytes, build: _Builder, job: _Job) -> str:
    """The 1990s Office formats, by way of LibreOffice if it's installed.

    No pure-Python reader for .doc or .ppt is worth relying on, and 400 MB of
    office suite in the image for two rare formats isn't a fair default. So it
    is used when present and explained when absent.
    """
    soffice = job.soffice or shutil.which("soffice") or shutil.which("libreoffice")
    modern = {".doc": "docx", ".ppt": "pptx", ".xls": "xlsx"}
    suffix = Path(job.filename).suffix.lower()
    target = modern.get(suffix, "docx")

    if not soffice:
        raise UnreadableFile(
            f"oneread can't read {suffix} files. Open it and save a copy as "
            f".{target}, then upload that."
        )

    with tempfile.TemporaryDirectory(prefix="oneread-convert-") as workspace:
        room = Path(workspace)
        source = room / f"upload{suffix}"
        source.write_bytes(data)
        try:
            subprocess.run(
                [
                    soffice,
                    f"-env:UserInstallation=file://{room / 'profile'}",
                    "--headless",
                    "--norestore",
                    "--convert-to",
                    target,
                    "--outdir",
                    str(room),
                    str(source),
                ],
                capture_output=True,
                timeout=SOFFICE_TIMEOUT_S,
                check=False,
            )
        except subprocess.TimeoutExpired:
            raise UnreadableFile("That file took too long to convert. Try a smaller one.") from None

        converted = room / f"upload.{target}"
        if not converted.is_file():
            raise UnreadableFile(
                f"That {suffix} file couldn't be converted. Save a copy as .{target} and "
                "upload that."
            )
        inner = converted.read_bytes()

    nested = _Job(
        filename=f"upload.{target}",
        limit=job.limit,
        max_unzipped=job.max_unzipped,
        soffice=job.soffice,
    )
    reader = READERS[f".{target}"]
    _guard_zip(inner, nested)
    return reader.read(inner, build, nested)


# --- the table of what we can read -------------------------------------------

_ZIP = b"PK\x03\x04"

READERS: dict[str, Reader] = {
    reader.ext: reader
    for reader in (
        Reader(".txt", "Plain text", "text", _read_text),
        Reader(".md", "Markdown", "markdown", _read_markdown, format=MARKDOWN),
        Reader(".markdown", "Markdown", "markdown", _read_markdown, format=MARKDOWN),
        Reader(".csv", "Spreadsheet export", "spreadsheet", _read_csv),
        Reader(".tsv", "Spreadsheet export", "spreadsheet", _read_csv),
        Reader(".docx", "Word", "document", _read_docx, zipped=True, magic=_ZIP),
        Reader(".pptx", "Slides", "slides", _read_pptx, zipped=True, magic=_ZIP),
        Reader(".xlsx", "Spreadsheet", "spreadsheet", _read_xlsx, zipped=True, magic=_ZIP),
        Reader(
            ".pdf",
            "PDF",
            "pdf",
            _read_pdf,
            magic=b"%PDF",
            empty="There's no text in that PDF to read.",
        ),
        Reader(".odt", "OpenDocument text", "document", _read_odf, zipped=True, magic=_ZIP),
        Reader(".ods", "OpenDocument sheet", "spreadsheet", _read_odf, zipped=True, magic=_ZIP),
        Reader(".odp", "OpenDocument slides", "slides", _read_odf, zipped=True, magic=_ZIP),
        Reader(".rtf", "Rich text", "document", _read_rtf),
        Reader(".html", "Web page", "document", _read_html),
        Reader(".htm", "Web page", "document", _read_html),
        Reader(".xls", "Older spreadsheet", "spreadsheet", _read_xls),
        Reader(".doc", "Older Word", "document", _read_legacy),
        Reader(".ppt", "Older slides", "slides", _read_legacy),
    )
}

#: What the file picker offers, in the order people are likely to want it.
ORDER = (
    ".txt", ".md", ".markdown", ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".xlsx", ".xls", ".csv", ".tsv", ".odt", ".odp", ".ods", ".rtf", ".html", ".htm",
)


def accepted() -> list[dict[str, str]]:
    """Every extension we read, for the picker and the copy beside it."""
    return [{"ext": ext, "label": READERS[ext].label, "kind": READERS[ext].kind} for ext in ORDER]


def media_type_for(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    return _MEDIA_TYPES.get(ext, "application/octet-stream")


_MEDIA_TYPES = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".csv": "text/csv",
    ".tsv": "text/tab-separated-values",
    ".html": "text/html",
    ".htm": "text/html",
    ".rtf": "application/rtf",
    ".pdf": "application/pdf",
    ".doc": "application/msword",
    ".ppt": "application/vnd.ms-powerpoint",
    ".xls": "application/vnd.ms-excel",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".odt": "application/vnd.oasis.opendocument.text",
    ".ods": "application/vnd.oasis.opendocument.spreadsheet",
    ".odp": "application/vnd.oasis.opendocument.presentation",
}


# --- the way in --------------------------------------------------------------


def extract(
    data: bytes,
    filename: str,
    *,
    limit: int = 100_000,
    max_unzipped: int = 200 * 1024 * 1024,
    soffice: str = "",
) -> Extracted:
    """Read `data` and return the words in it, or say why we can't."""
    ext = Path(filename).suffix.lower()
    reader = READERS.get(ext)
    if reader is None:
        raise UnreadableFile(
            f"oneread can't read {ext or 'files without an extension'}. "
            "Word, slides, spreadsheets, PDF, markdown and plain text all work."
        )
    if not data.strip():
        raise UnreadableFile("That file is empty.")
    if reader.magic and not data.startswith(reader.magic):
        raise UnreadableFile(
            f"That doesn't look like a {reader.label} file inside, whatever the name says."
        )

    job = _Job(filename=filename, limit=limit, max_unzipped=max_unzipped, soffice=soffice)
    if reader.zipped:
        _guard_zip(data, job)

    build = _Builder(limit)
    try:
        found = reader.read(data, build, job)
    except UnreadableFile:
        raise
    except Exception:
        # The person uploading gets a sentence; the log gets the traceback,
        # because "it may be damaged" is useless when the fault is ours.
        log.warning("could not read %s", filename, exc_info=True)
        raise UnreadableFile(
            f"That {reader.label} file couldn't be read. It may be damaged."
        ) from None

    text = build.done()
    if not text.strip():
        raise UnreadableFile(reader.empty)

    return Extracted(
        text=text,
        format=reader.format,
        title=_title(found, filename),
        kind=reader.kind,
        truncated=build.truncated,
    )


_SEPARATORS = re.compile(r"[_\-]+")


def _title(found: str, filename: str) -> str:
    """The document's own title, or a tidied-up version of its filename."""
    title = _tidy(found).replace("\n", " ").strip()
    if not title:
        title = _SEPARATORS.sub(" ", Path(filename).stem).strip()
    return _SPACE.sub(" ", title)[:200] or "Untitled"
