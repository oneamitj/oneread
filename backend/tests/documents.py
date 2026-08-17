"""Sample documents, built in memory.

Every one of these is made with the same library that reads it back, so the
suite carries no binary files and nothing goes stale when a dependency moves.
The PDF is the exception — pypdf writes PDFs but won't put text in one, so
that one is assembled by hand, xref table and all.
"""

from __future__ import annotations

import io
import zipfile


def docx(*, title: str = "Quarterly Report") -> bytes:
    from docx import Document

    document = Document()
    document.core_properties.title = title
    document.add_heading(title, level=1)
    document.add_paragraph("Revenue grew in both regions.")
    table = document.add_table(rows=3, cols=3)
    rows = (
        ("Region", "Units", "Revenue"),
        ("North", "412", "88300"),
        ("South", "388", "71100"),
    )
    for row, values in enumerate(rows):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    document.add_paragraph("That is the whole picture.")
    return _saved(document)


def pptx() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "How it went"
    slide.placeholders[1].text = "Up and to the right"
    slide.notes_slide.notes_text_frame.text = "Mention the caveat about March."
    return _saved(presentation)


def xlsx() -> bytes:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Q3"
    for row in (("Region", "Units", "Revenue"), ("North", 412, 88300.0)):
        sheet.append(row)
    later = book.create_sheet("Q4")
    later.append(("Region", "Units"))
    later.append(("South", 500))
    return _saved(book)


def odt() -> bytes:
    from odf import table as odf_table
    from odf import text as odf_text
    from odf.opendocument import OpenDocumentText

    document = OpenDocumentText()
    document.text.addElement(odf_text.H(outlinelevel=1, text="Open Report"))
    document.text.addElement(odf_text.P(text="Some prose here."))
    document.text.addElement(_odf_grid(odf_table, odf_text))
    return _saved(document)


def ods() -> bytes:
    from odf import table as odf_table
    from odf import text as odf_text
    from odf.opendocument import OpenDocumentSpreadsheet

    document = OpenDocumentSpreadsheet()
    document.spreadsheet.addElement(_odf_grid(odf_table, odf_text))
    return _saved(document)


def odp() -> bytes:
    from odf import draw, style
    from odf import text as odf_text
    from odf.opendocument import OpenDocumentPresentation

    document = OpenDocumentPresentation()
    document.masterstyles.addElement(style.MasterPage(name="M", pagelayoutname="PL"))
    page = draw.Page(name="page1", masterpagename="M")
    frame = draw.Frame(width="10cm", height="2cm", x="1cm", y="1cm")
    box = draw.TextBox()
    box.addElement(odf_text.P(text="Hello slide"))
    frame.addElement(box)
    page.addElement(frame)
    document.presentation.addElement(page)
    return _saved(document)


def _odf_grid(odf_table, odf_text):
    grid = odf_table.Table(name="Grid")
    for values in (["Region", "Units"], ["North", "412"]):
        row = odf_table.TableRow()
        for value in values:
            cell = odf_table.TableCell(valuetype="string")
            cell.addElement(odf_text.P(text=value))
            row.addElement(cell)
        grid.addElement(row)
    return grid


def _saved(document) -> bytes:
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def pdf(lines: list[str]) -> bytes:
    """A one-page PDF. An empty `lines` gives a page with no text on it at all."""
    if lines:
        drawn = " ".join(f"({line}) Tj 0 -16 Td" for line in lines)
        stream = f"BT /F1 12 Tf 72 720 Td {drawn} ET".encode()
    else:
        stream = b"0 0 0 rg 10 10 100 100 re f"

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % number + body + b"\nendobj\n"

    start = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objects) + 1)
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objects) + 1,
        start,
    )
    return bytes(out)


def zip_bomb(size: int = 60 * 1024 * 1024) -> bytes:
    """A tiny archive that unpacks to something enormous."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as bundle:
        bundle.writestr("word/document.xml", b"a" * size)
    return buffer.getvalue()


def entity_bomb() -> bytes:
    """An office file whose XML declares entities, the billion-laughs shape."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as bundle:
        bundle.writestr(
            "word/document.xml",
            b'<?xml version="1.0"?><!DOCTYPE lol [<!ENTITY a "ha ha ha">]>'
            b"<w:document>&a;</w:document>",
        )
    return buffer.getvalue()
